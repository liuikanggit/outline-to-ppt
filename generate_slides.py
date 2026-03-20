#!/usr/bin/env python3
"""
根据 input.md 大纲和已生成的母版模板，生成完整的 PPT 幻灯片。

幻灯片结构：
1. 封面页（保持不变）
2. 总目录页（根据大纲重新生成）
3. 对每个章节：
   a. 章节页（显示章节名称）
   b. 对每个小节：内容页（关联对应母版，标题=小节名，内容=xxxx）
4. 结束页（保持不变）
"""

import json
import os
import re
import copy
import urllib.request

def download_image(uri, output_dir='images'):
    """下载幕布图片并缓存到本地 (使用 urllib 兼容无 requests 环境)"""
    if not uri:
        return None
    import os
    os.makedirs(output_dir, exist_ok=True)
    
    filename = os.path.basename(uri)
    local_path = os.path.join(output_dir, filename)
    
    # 如果已存在，直接返回
    if os.path.exists(local_path):
        return local_path
        
    url = f"https://api2.mubu.com/v3/{uri}"
    try:
        req = urllib.request.Request(url)
        with urllib.request.urlopen(req, timeout=5) as resp:
            with open(local_path, 'wb') as f:
                f.write(resp.read())
            return local_path
    except Exception as e:
        print(f"下载图片失败 {url}: {e}")
        
    return None

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ==================== 专属动态导航条常量与命名空间 ====================
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

CHAPTER_BOX_Y = 19172
CHAPTER_BOX_WIDTH = 1163782
CHAPTER_BOX_HEIGHT = 369332
CHAPTER_BOX_STEP = 1062188
RIGHT_MARGIN = 12116572

SECTION_BOX_Y = 427063
SECTION_BOX_HEIGHT = 259675
SECTION_RIGHT_MARGIN = 12049637  # 节行右边缘

SUBSECTION_BOX_Y = 728980
SUBSECTION_BOX_HEIGHT = 243840
SUBSECTION_RIGHT_MARGIN = 12049760  # 小节行右边缘

DECORATOR_GROUP_WIDTH = 997398
DECORATOR_GROUP_HEIGHT = 400620
DECORATOR_OFFSET_X = 84562

# ---- 辅助画条颜色/XML函数 (由 generate_masters.py 提纯) ----
def make_solid_fill_highlight():
    a = NSMAP['a']
    solidFill = etree.Element(f'{{{a}}}solidFill')
    srgbClr = etree.SubElement(solidFill, f'{{{a}}}srgbClr', val='0096FF')
    return solidFill

def make_solid_fill_normal():
    a = NSMAP['a']
    solidFill = etree.Element(f'{{{a}}}solidFill')
    srgbClr = etree.SubElement(solidFill, f'{{{a}}}srgbClr', val='7F7F7F')
    return solidFill

def make_solid_fill_chapter_highlight():
    a = NSMAP['a']
    solidFill = etree.Element(f'{{{a}}}solidFill')
    schemeClr = etree.SubElement(solidFill, f'{{{a}}}schemeClr', val='accent1')
    return solidFill

def make_solid_fill_black():
    """提供纯黑色填充用于高亮文字"""
    a = NSMAP['a']
    solidFill = etree.Element(f'{{{a}}}solidFill')
    etree.SubElement(solidFill, f'{{{a}}}srgbClr', val='000000')
    return solidFill

def make_solid_fill_chapter_normal():
    a = NSMAP['a']
    solidFill = etree.Element(f'{{{a}}}solidFill')
    schemeClr = etree.SubElement(solidFill, f'{{{a}}}schemeClr', val='bg2')
    lumMod = etree.SubElement(schemeClr, f'{{{a}}}lumMod', val='60000')
    return solidFill

def set_run_fill(rpr_elem, fill_elem):
    for child in rpr_elem:
        if child.tag.endswith('solidFill') or child.tag.endswith('noFill'):
            rpr_elem.remove(child)
    rpr_elem.append(fill_elem)

def calculate_chapter_positions(num_chapters, use_fixed_template=False):
    if use_fixed_template:
        # 内容页绝对刻度：精确对齐母版 1 中的 [文本框 12] (X=3517476) 起步，步长 1062188
        start_x = 3517476
        step_x = 1062188
        return [start_x + i * step_x for i in range(num_chapters)]
    else:
        # 总目录页等其他页面：保持居中排布
        total_width = num_chapters * CHAPTER_BOX_WIDTH + (num_chapters - 1) * (CHAPTER_BOX_STEP - CHAPTER_BOX_WIDTH)
        avail_width = 12192000
        start_x = (avail_width - total_width) // 2
        return [start_x + i * CHAPTER_BOX_STEP for i in range(num_chapters)]

def create_chapter_textbox(chapter_idx, chapter_name, is_highlight, x_pos, sp_id):
    a, p = NSMAP['a'], NSMAP['p']
    sp = etree.Element(f'{{{p}}}sp')
    nvSpPr = etree.SubElement(sp, f'{{{p}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{p}}}cNvPr', id=str(sp_id), name=f'导航_章_{sp_id}')
    etree.SubElement(nvSpPr, f'{{{p}}}cNvSpPr', txBox='1')
    etree.SubElement(nvSpPr, f'{{{p}}}nvPr', userDrawn='1')
    
    spPr = etree.SubElement(sp, f'{{{p}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{a}}}xfrm')
    etree.SubElement(xfrm, f'{{{a}}}off', x=str(x_pos), y=str(CHAPTER_BOX_Y))
    etree.SubElement(xfrm, f'{{{a}}}ext', cx=str(CHAPTER_BOX_WIDTH), cy=str(CHAPTER_BOX_HEIGHT))
    etree.SubElement(etree.SubElement(spPr, f'{{{a}}}prstGeom', prst='rect'), f'{{{a}}}avLst')
    etree.SubElement(spPr, f'{{{a}}}noFill')

    txBody = etree.SubElement(sp, f'{{{p}}}txBody')
    etree.SubElement(txBody, f'{{{a}}}bodyPr', wrap='none', rtlCol='0')
    etree.SubElement(txBody, f'{{{a}}}lstStyle')
    
    p_elem = etree.SubElement(txBody, f'{{{a}}}p')
    etree.SubElement(p_elem, f'{{{a}}}pPr', algn='ctr')
    r_elem = etree.SubElement(p_elem, f'{{{a}}}r')
    rPr = etree.SubElement(r_elem, f'{{{a}}}rPr', lang='zh-CN', sz='1200', b='1' if is_highlight else '0')
    etree.SubElement(rPr, f'{{{a}}}latin', typeface='微软雅黑')
    etree.SubElement(rPr, f'{{{a}}}ea', typeface='微软雅黑')
    
    set_run_fill(rPr, make_solid_fill_black() if is_highlight else make_solid_fill_chapter_normal())
    etree.SubElement(r_elem, f'{{{a}}}t').text = chapter_name
    return sp

def create_section_bar(sections, highlight_section_idx, x_pos, width, sp_id):
    a = NSMAP['a']
    p_ns = NSMAP['p']
    sp = etree.Element(f'{{{p_ns}}}sp')
    nvSpPr = etree.SubElement(sp, f'{{{p_ns}}}nvSpPr')
    etree.SubElement(nvSpPr, f'{{{p_ns}}}cNvPr', id=str(sp_id), name=f'导航_节_{sp_id}')
    etree.SubElement(nvSpPr, f'{{{p_ns}}}cNvSpPr', txBox='1')
    etree.SubElement(nvSpPr, f'{{{p_ns}}}nvPr', userDrawn='1')

    spPr = etree.SubElement(sp, f'{{{p_ns}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{a}}}xfrm')
    etree.SubElement(xfrm, f'{{{a}}}off', x=str(x_pos), y=str(SECTION_BOX_Y))
    etree.SubElement(xfrm, f'{{{a}}}ext', cx=str(width), cy=str(SECTION_BOX_HEIGHT))
    prstGeom = etree.SubElement(spPr, f'{{{a}}}prstGeom', prst='roundRect')
    avLst = etree.SubElement(prstGeom, f'{{{a}}}avLst')
    etree.SubElement(avLst, f'{{{a}}}gd', name='adj', fmla='val 50000')
    etree.SubElement(etree.SubElement(spPr, f'{{{a}}}solidFill'), f'{{{a}}}srgbClr', val='F5F5F9')
    ln = etree.SubElement(spPr, f'{{{a}}}ln', w='3175')
    lnClr = etree.SubElement(etree.SubElement(ln, f'{{{a}}}solidFill'), f'{{{a}}}schemeClr', val='bg1')
    etree.SubElement(lnClr, f'{{{a}}}lumMod', val='65000')

    txBody = etree.SubElement(sp, f'{{{p_ns}}}txBody')
    etree.SubElement(etree.SubElement(txBody, f'{{{a}}}bodyPr', wrap='square'), f'{{{a}}}spAutoFit')
    etree.SubElement(txBody, f'{{{a}}}lstStyle')

    para = etree.SubElement(txBody, f'{{{a}}}p')
    etree.SubElement(para, f'{{{a}}}pPr', algn='ctr')
    separator = '      '

    for i, sec in enumerate(sections):
        is_current = (i == highlight_section_idx)
        # 获取二级菜单名称
        sec_name = sec.get('chapterName', f'第{i+1}节')
        
        fill = make_solid_fill_highlight() if is_current else make_solid_fill_normal()
        font_name = '微软雅黑' if is_current else 'Microsoft YaHei Light'
        bold = '1' if is_current else '0'

        if i > 0:
            sep_r = etree.SubElement(para, f'{{{a}}}r')
            sep_rPr = etree.SubElement(sep_r, f'{{{a}}}rPr', sz='1200')
            sep_rPr.append(copy.deepcopy(fill))
            etree.SubElement(sep_rPr, f'{{{a}}}latin', typeface=font_name)
            etree.SubElement(sep_r, f'{{{a}}}t').text = separator

        r = etree.SubElement(para, f'{{{a}}}r')
        rPr = etree.SubElement(r, f'{{{a}}}rPr', sz='1200', b=bold)
        rPr.append(copy.deepcopy(fill))
        etree.SubElement(rPr, f'{{{a}}}latin', typeface=font_name)
        etree.SubElement(rPr, f'{{{a}}}ea', typeface=font_name)
        etree.SubElement(r, f'{{{a}}}t').text = sec_name

    return sp

def create_subsection_bar(subsections, highlight_subsection_idx, x_pos, width, sp_id):
    a = NSMAP['a']
    p_ns = NSMAP['p']
    sp = etree.Element(f'{{{p_ns}}}sp')
    nvSpPr = etree.SubElement(sp, f'{{{p_ns}}}nvSpPr')
    etree.SubElement(nvSpPr, f'{{{p_ns}}}cNvPr', id=str(sp_id), name=f'导航_小节_{sp_id}')
    etree.SubElement(nvSpPr, f'{{{p_ns}}}cNvSpPr', txBox='1')
    
    spPr = etree.SubElement(sp, f'{{{p_ns}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{a}}}xfrm')
    etree.SubElement(xfrm, f'{{{a}}}off', x=str(x_pos), y=str(SUBSECTION_BOX_Y))
    etree.SubElement(xfrm, f'{{{a}}}ext', cx=str(width), cy=str(SUBSECTION_BOX_HEIGHT))
    prstGeom = etree.SubElement(spPr, f'{{{a}}}prstGeom', prst='roundRect')
    avLst = etree.SubElement(prstGeom, f'{{{a}}}avLst')
    etree.SubElement(avLst, f'{{{a}}}gd', name='adj', fmla='val 50000')
    etree.SubElement(etree.SubElement(spPr, f'{{{a}}}solidFill'), f'{{{a}}}srgbClr', val='F5F5F9')

    txBody = etree.SubElement(sp, f'{{{p_ns}}}txBody')
    etree.SubElement(txBody, f'{{{a}}}bodyPr', wrap='square')
    para = etree.SubElement(txBody, f'{{{a}}}p')
    etree.SubElement(para, f'{{{a}}}pPr', algn='ctr')
    separator = '      '

    for i, sub in enumerate(subsections):
        is_current = (i == highlight_subsection_idx)
        sub_name = sub.get('chapterName', f'第{i+1}小节')

        fill = make_solid_fill_highlight() if is_current else make_solid_fill_normal()
        font_name = '微软雅黑' if is_current else 'Microsoft YaHei Light'

        if i > 0:
            sep_r = etree.SubElement(para, f'{{{a}}}r')
            sep_rPr = etree.SubElement(sep_r, f'{{{a}}}rPr', sz='1200')
            sep_rPr.append(copy.deepcopy(fill))
            etree.SubElement(sep_r, f'{{{a}}}t').text = separator

        r = etree.SubElement(para, f'{{{a}}}r')
        rPr = etree.SubElement(r, f'{{{a}}}rPr', sz='1200', b='1' if is_current else '0')
        rPr.append(copy.deepcopy(fill))
        etree.SubElement(rPr, f'{{{a}}}latin', typeface=font_name)
        etree.SubElement(r, f'{{{a}}}t').text = sub_name

    return sp

def calculate_section_bar_params(num_sections):
    per_section = 900000  # 适当加宽以避开重叠
    padding = 150000
    bar_width = num_sections * per_section + padding
    bar_x = SECTION_RIGHT_MARGIN - bar_width
    return bar_x, bar_width

def calculate_subsection_bar_params(num_subsections):
    per_subsection = 950000
    padding = 150000
    bar_width = num_subsections * per_subsection + padding
    bar_x = SUBSECTION_RIGHT_MARGIN - bar_width
    return bar_x, bar_width

def add_nav_bar_to_slide(slide, code_str, all_chapters, prs, sp_id_start=3000):
    """向单页 Slide 动态追加专属的 导航条 浮层形状（含黑色背景）"""
    parts = code_str.split('.')
    try: ch_idx = int(parts[0]) - 1
    except: ch_idx = 0
    try: sec_idx = int(parts[1]) - 1 if len(parts) > 1 else 0
    except: sec_idx = 0
    try: sub_idx = int(parts[2]) - 1 if len(parts) > 2 else 0
    except: sub_idx = 0

    spTree = slide.shapes._spTree
    sp_id = sp_id_start
    num_chapters = len(all_chapters)

    # ---- 🟢 1. 大扫除：将所有非 placeholders 且不相关的 Layout 继承物全量清空 🟢 ----
    shapes_to_remove = []
    page_ph_names = [ph.name for ph in slide.placeholders]
    for shp in list(slide.shapes):
         if shp.name not in page_ph_names and '标题' not in shp.name and '内容' not in shp.name:
              shapes_to_remove.append(shp)
    for shp in shapes_to_remove:
         try: spTree.remove(shp._element)
         except: pass

    # ---- 0. 高亮背景 组合 26 仿制 (深拷贝母版1) ----
    decorator_grp = None
    logo_grp = None
    try:
        source_master = prs.slide_masters[1]
        for s in source_master.shapes:
            if s.name == '组合 26':
                 decorator_grp = s
                 break
    except:
        pass

    if logo_grp is not None:
         spTree.append(copy.deepcopy(logo_grp._element))

    positions = calculate_chapter_positions(num_chapters, use_fixed_template=True) if num_chapters > 0 else []

    if decorator_grp is not None and num_chapters > 0:
        dec_xml = copy.deepcopy(decorator_grp._element)
        xfrm = dec_xml.find('.//' + f'{{{NSMAP["a"]}}}xfrm')
        if xfrm is not None:
             off = xfrm.find(f'{{{NSMAP["a"]}}}off')
             if off is not None:
                  # 模板原始：组合宽=997398， 我们的TextBox宽=1163782
                  dec_x = positions[ch_idx] + (1163782 - 997398) // 2
                  off.set('x', str(dec_x))
                  off.set('y', '0')
        spTree.append(dec_xml)  # 优先追加作为底层

    # ---- 1. 添加 Level 1 (章节) 导航条 ----
    if num_chapters > 0:
        for i, ch in enumerate(all_chapters):
            ch_name = ch.get('chapterName', '')
            is_hl = (i == ch_idx)
            sp = create_chapter_textbox(i, ch_name, is_hl, positions[i], sp_id)
            spTree.append(sp)
            sp_id += 1

    # ---- 2. 添加 Level 2 (节) 导航条 ----
    current_chapter = all_chapters[ch_idx] if ch_idx < len(all_chapters) else {}
    sub_chapters = current_chapter.get('subChapter', [])
    num_sections = len(sub_chapters)
    
    if num_sections > 0:
        bar_x, bar_w = calculate_section_bar_params(num_sections)
        sp_sec = create_section_bar(sub_chapters, sec_idx, bar_x, bar_w, sp_id)
        spTree.append(sp_sec)
        sp_id += 1

    # ---- 3. 添加 Level 3 (小节) 导航条 ----
    if sec_idx < len(sub_chapters):
        current_section = sub_chapters[sec_idx]
        subsub_chapters = current_section.get('subChapter', [])
        num_sub = len(subsub_chapters)
        if num_sub > 0:
             bar_sub_x, bar_sub_w = calculate_subsection_bar_params(num_sub)
             sp_sub = create_subsection_bar(subsub_chapters, sub_idx, bar_sub_x, bar_sub_w, sp_id)
             spTree.append(sp_sub)
             sp_id += 1

# ==================== 中文数字映射 ====================
CN_NUMBERS = ['一', '二', '三', '四', '五', '六', '七', '八', '九', '十',
              '十一', '十二', '十三', '十四', '十五']

EN_NUMBERS = ['One', 'Two', 'Three', 'Four', 'Five', 'Six', 'Seven', 'Eight',
              'Nine', 'Ten', 'Eleven', 'Twelve', 'Thirteen', 'Fourteen', 'Fifteen']


# ==================== 解析 input.md ====================
def parse_input(filepath):
    """通过 mubu_parser 加载幕布数据"""
    from mubu_parser import load_mubu_data
    return load_mubu_data(filepath)


# ==================== 总目录页生成 ====================
def create_toc_slide(prs, slide, outline):
    """修改总目录页（第2页），根据大纲重新生成内容"""
    # 保留 "总目录" 标题（文本框7），删除其他所有形状
    shapes_to_keep = []
    title_shape = None

    for shape in slide.shapes:
        if shape.name == '文本框 7':
            title_shape = shape
            shapes_to_keep.append(shape)

    # 通过 XML 删除不需要的形状
    spTree = slide.shapes._spTree
    shapes_to_remove = []
    for sp in list(spTree):
        tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
        if tag == 'sp':
            # 检查 cNvPr 的 name
            cNvPr = sp.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}cNvPr')
            if cNvPr is None:
                cNvPr = sp.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
            if cNvPr is None:
                cNvPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
            if cNvPr is None:
                # 直接搜索所有带 name 的 cNvPr
                for elem in sp.iter():
                    if elem.tag.endswith('cNvPr') and 'name' in elem.attrib:
                        cNvPr = elem
                        break

            if cNvPr is not None:
                name = cNvPr.get('name', '')
                if name != '文本框 7':
                    shapes_to_remove.append(sp)

    for sp in shapes_to_remove:
        spTree.remove(sp)

    # 现在添加新的章节标题和节列表
    num_chapters = len(outline)

    # 总目录区域参数（固定布局和间距，整体居中）
    chapter_box_width = 1300000  # 略微拉宽，防单行溢出，由 1163782 提升
    chapter_box_height = 461665
    chapter_gap = 120000         # 间距收敛，配合拉宽
    toc_area_top = (prs.slide_height - chapter_box_height) // 2   # 垂直居中
    section_list_top = 2205146  # 已无用

    # 计算整体居中的左起始点
    slide_width = prs.slide_width
    total_width = num_chapters * chapter_box_width + (num_chapters - 1) * chapter_gap
    toc_area_left = (slide_width - total_width) // 2

    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    sp_id = 100
    for ch_idx, chapter in enumerate(outline):
        x = toc_area_left + ch_idx * (chapter_box_width + chapter_gap)

        # ---- 章节标题文本框 ----
        chapter_sp = create_toc_chapter_box(
            chapter, ch_idx, x, toc_area_top,
            chapter_box_width, chapter_box_height, sp_id)
        spTree.append(chapter_sp)
        sp_id += 1


def create_toc_chapter_box(chapter_dict, ch_idx, x, y, w, h, sp_id):
    """创建总目录中的章节标题文本框"""
    a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

    sp = etree.SubElement(etree.Element('dummy'), f'{{{p}}}sp')
    sp = sp  # detach

    # nvSpPr
    nvSpPr = etree.SubElement(sp, f'{{{p}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{p}}}cNvPr')
    cNvPr.set('id', str(sp_id))
    cNvPr.set('name', f'TOC_Chapter_{ch_idx}')
    cNvSpPr = etree.SubElement(nvSpPr, f'{{{p}}}cNvSpPr')
    cNvSpPr.set('txBox', '1')
    etree.SubElement(nvSpPr, f'{{{p}}}nvPr')

    # spPr
    spPr = etree.SubElement(sp, f'{{{p}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{a}}}xfrm')
    etree.SubElement(xfrm, f'{{{a}}}off', x=str(x), y=str(y))
    etree.SubElement(xfrm, f'{{{a}}}ext', cx=str(w), cy=str(h))
    prstGeom = etree.SubElement(spPr, f'{{{a}}}prstGeom', prst='rect')
    etree.SubElement(prstGeom, f'{{{a}}}avLst')

    # txBody
    txBody = etree.SubElement(sp, f'{{{p}}}txBody')
    bodyPr = etree.SubElement(txBody, f'{{{a}}}bodyPr', wrap='none', rtlCol='0')
    etree.SubElement(txBody, f'{{{a}}}lstStyle')

    # 第一行：中文章节名（16pt 粗体 微软雅黑）
    para1 = etree.SubElement(txBody, f'{{{a}}}p')
    pPr1 = etree.SubElement(para1, f'{{{a}}}pPr', algn='ctr')
    run1 = etree.SubElement(para1, f'{{{a}}}r')
    rPr1 = etree.SubElement(run1, f'{{{a}}}rPr', lang='zh-CN', sz='1600', b='1', dirty='0')
    solidFill1 = etree.SubElement(rPr1, f'{{{a}}}solidFill')
    schemeClr1 = etree.SubElement(solidFill1, f'{{{a}}}schemeClr', val='tx1')
    etree.SubElement(rPr1, f'{{{a}}}latin', typeface='微软雅黑')
    run1_t = etree.SubElement(run1, f'{{{a}}}t')
    run1_t.text = chapter_dict['chapterName']

    # 第二行：英文章节名（8pt Microsoft YaHei Light）
    name_en = chapter_dict.get('chapterEnName', '')
    if name_en:
        para2 = etree.SubElement(txBody, f'{{{a}}}p')
        etree.SubElement(para2, f'{{{a}}}pPr', algn='ctr')
        run2 = etree.SubElement(para2, f'{{{a}}}r')
        rPr2 = etree.SubElement(run2, f'{{{a}}}rPr', lang='en-US', sz='800', b='0', dirty='0')
        solidFill2 = etree.SubElement(rPr2, f'{{{a}}}solidFill')
        etree.SubElement(solidFill2, f'{{{a}}}schemeClr', val='tx1')
        etree.SubElement(rPr2, f'{{{a}}}latin', typeface='Microsoft YaHei Light')
        etree.SubElement(rPr2, f'{{{a}}}ea', typeface='Microsoft YaHei Light')
        run2_t = etree.SubElement(run2, f'{{{a}}}t')
        run2_t.text = name_en

    return sp


def create_toc_section_list(sections, x, y, w, h, sp_id):
    """创建总目录中的节列表文本框"""
    a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    sp = etree.SubElement(etree.Element('dummy'), f'{{{p}}}sp')

    nvSpPr = etree.SubElement(sp, f'{{{p}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{p}}}cNvPr')
    cNvPr.set('id', str(sp_id))
    cNvPr.set('name', f'TOC_Sections_{sp_id}')
    cNvSpPr = etree.SubElement(nvSpPr, f'{{{p}}}cNvSpPr')
    cNvSpPr.set('txBox', '1')
    etree.SubElement(nvSpPr, f'{{{p}}}nvPr')

    spPr = etree.SubElement(sp, f'{{{p}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{a}}}xfrm')
    etree.SubElement(xfrm, f'{{{a}}}off', x=str(x), y=str(y))
    etree.SubElement(xfrm, f'{{{a}}}ext', cx=str(w), cy=str(h))
    prstGeom = etree.SubElement(spPr, f'{{{a}}}prstGeom', prst='rect')
    etree.SubElement(prstGeom, f'{{{a}}}avLst')

    txBody = etree.SubElement(sp, f'{{{p}}}txBody')
    bodyPr = etree.SubElement(txBody, f'{{{a}}}bodyPr', wrap='none', rtlCol='0')
    etree.SubElement(txBody, f'{{{a}}}lstStyle')

    # 每个节一行（14pt Microsoft YaHei Light）
    for sec in sections:
        para = etree.SubElement(txBody, f'{{{a}}}p')
        pPr = etree.SubElement(para, f'{{{a}}}pPr', algn='ctr')
        run = etree.SubElement(para, f'{{{a}}}r')
        rPr = etree.SubElement(run, f'{{{a}}}rPr', lang='zh-CN', sz='1400', dirty='0')
        solidFill = etree.SubElement(rPr, f'{{{a}}}solidFill')
        schemeClr = etree.SubElement(solidFill, f'{{{a}}}schemeClr', val='tx1')
        lumMod = etree.SubElement(schemeClr, f'{{{a}}}lumMod', val='75000')
        lumOff = etree.SubElement(schemeClr, f'{{{a}}}lumOff', val='25000')
        etree.SubElement(rPr, f'{{{a}}}latin', typeface='Microsoft YaHei Light')
        etree.SubElement(rPr, f'{{{a}}}ea', typeface='Microsoft YaHei Light')
        run_t = etree.SubElement(run, f'{{{a}}}t')
        run_t.text = sec['chapterName']

    return sp


# ==================== 章节页生成 ====================
def create_chapter_slide(prs, template_slide, chapter_dict, ch_idx):
    """基于第3页模板创建新的章节页。直接修改复制的元素以保持原始格式"""
    layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # 清除新页上的占位符
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    spTree_new = new_slide.shapes._spTree
    spTree_old = template_slide.shapes._spTree
    
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    
    for sp in spTree_old:
        tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
        if tag == 'sp':
            # 只复制形状
            new_sp = copy.deepcopy(sp)
            spTree_new.append(new_sp)
            
            # 修改文本内容
            txBody = new_sp.find(f'.//{{http://schemas.openxmlformats.org/presentationml/2006/main}}txBody')
            if txBody is not None:
                paras = txBody.findall(f'.//{{{a_ns}}}p')

                if len(paras) >= 1:
                    # 第一段：中文名
                    runs = paras[0].findall(f'.//{{{a_ns}}}r/{{{a_ns}}}t')
                    if runs:
                        runs[0].text = chapter_dict.get('chapterName', '')
                if len(paras) >= 2:
                    # 第二段：英文名
                    runs = paras[1].findall(f'.//{{{a_ns}}}r/{{{a_ns}}}t')
                    if runs:
                        runs[0].text = chapter_dict.get('chapterEnName', '')

    return new_slide


# ==================== 内容页生成 ====================
def estimate_lines(text, font_size=16):
    """根据字号估算单文本行数"""
    if not text:
        return 0
    char_per_line = 38 if font_size == 16 else 44  # 估算一排可放字符
    import math
    return math.ceil(len(text) / char_per_line)

def paginate_contents(contents):
    """分页与字号决策算法"""
    MAX_LINES = 12  # 从 14 降到 12
    # 1. 尝试 16pt (1600)
    pages_16 = []
    current_page = []
    current_lines = 0
    
    for item in contents:
        txt = item.get('text', '')
        lines = estimate_lines(txt, 16)
        if current_lines + lines <= MAX_LINES:
            current_page.append(item)
            current_lines += (lines if lines > 0 else 1)
        else:
            if current_page:
                pages_16.append(current_page)
            current_page = [item]
            current_lines = lines if lines > 0 else 1
    if current_page:
        pages_16.append(current_page)
        
    # 2. 如果 16pt 产生了多页，尝试 14pt 能否合并成更少页
    if len(pages_16) > 1:
        pages_14 = []
        current_page_14 = []
        current_lines_14 = 0
        for item in contents:
            txt = item.get('text', '')
            lines = estimate_lines(txt, 14)
            if current_lines_14 + lines <= MAX_LINES:
                current_page_14.append(item)
                current_lines_14 += (lines if lines > 0 else 1)
            else:
                if current_page_14:
                    pages_14.append(current_page_14)
                current_page_14 = [item]
                current_lines_14 = lines if lines > 0 else 1
        if current_page_14:
            pages_14.append(current_page_14)
            
        if len(pages_14) < len(pages_16):
            return pages_14, 1400
            
    return pages_16, 1600

def create_content_slide_v2(prs, content_layout, page_title, contents, default_layout, code_str, all_chapters):
    """新版内容页生成（适配 Content[] 模型）"""
    
    DEBUG_LAYOUT_ONLY = False
    if DEBUG_LAYOUT_ONLY:
         new_slide = prs.slides.add_slide(content_layout)
         # 追加专属导航条 (含 prs)
         add_nav_bar_to_slide(new_slide, code_str, all_chapters, prs)
         
         # 🟢 听从指令：页面中央大标题彻底清空，达成“PPT空着”的要求
         if hasattr(new_slide.shapes, 'title') and new_slide.shapes.title:
              new_slide.shapes.title.text = ""
         else:
              for shp in new_slide.shapes:
                   if hasattr(shp, 'has_text_frame') and shp.has_text_frame:
                        if '标题' in shp.text_frame.text:
                             shp.text_frame.text = ""
                             break
         return [new_slide]

    def flatten_content(content_item):
        """将树型 TextContent 扁平化为带 level 级联的旧版结构项"""
        flat = []
        if content_item.get('type') == 'text':
            text_info = content_item.get('text', {})
            main_text = text_info.get('text', '')
            if main_text:
                # 顶级文本 level=0 (不带项目符号)
                flat.append({ "text": main_text, "level": 0 })
                def recurse_sub(items, depth):
                    for sub in items:
                        flat.append({ "text": sub.get('text', ''), "level": depth })
                        if sub.get('subContent'):
                            recurse_sub(sub['subContent'], depth + 1)
                recurse_sub(text_info.get('subContent', []), 1)
        return flat

    # 1. 搜集并分类内容
    flat_text_items = []
    all_images = []
    
    for c in contents:
        if c.get('type') == 'image':
            img_info = c.get('image', {})
            if img_info.get('uri'):
                all_images.append(img_info)
        elif c.get('type') == 'text':
            flat_text_items.extend(flatten_content(c))

    generated_slides = []
    a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    # 2. ---- A. 生成文字页 ----
    if flat_text_items:
        pages, font_sz_val = paginate_contents(flat_text_items)
        for page_idx, page_items in enumerate(pages):
            new_slide = prs.slides.add_slide(content_layout)
            generated_slides.append(new_slide)

            # 标题填充容错: 依照指令，已有页眉，内容页大标题设为空 ""
            if hasattr(new_slide.shapes, 'title') and new_slide.shapes.title:
                new_slide.shapes.title.text = page_title
            else:
                for shp in new_slide.shapes:
                    if hasattr(shp, 'has_text_frame') and shp.has_text_frame:
                        if '标题' in shp.text_frame.text or shp.name == '标题':
                            shp.text_frame.text = page_title
                            break

            spTree = new_slide.shapes._spTree
            body_ph = None
            for ph in new_slide.placeholders:
                if ph.placeholder_format.type == 2:  # Body
                    body_ph = ph
                    break
            if not body_ph and len(new_slide.placeholders) > 1:
                body_ph = new_slide.placeholders[1]
            if body_ph is None:
                body_ph = new_slide.shapes[1] if len(new_slide.shapes) > 1 else None

            txBody2 = None
            if body_ph is not None:
                txBody2 = body_ph._element.find('.//' + f'{{{p_ns}}}txBody')
                if txBody2 is None:
                    txBody2 = etree.SubElement(body_ph._element, f'{{{p_ns}}}txBody')
                else:
                    for p in txBody2.findall(f'.//{{{a}}}p'):
                        txBody2.remove(p)

            # 最终降级占位：如果页面实在没有任何框（比如纯背景），手画一个接入正文，以绝 Crash
            if txBody2 is None:
                content_sp = etree.SubElement(spTree, f'{{{p_ns}}}sp')
                nvSpPr2 = etree.SubElement(content_sp, f'{{{p_ns}}}nvSpPr')
                etree.SubElement(nvSpPr2, f'{{{p_ns}}}cNvPr', id='2026', name='内容兜底')
                etree.SubElement(nvSpPr2, f'{{{p_ns}}}cNvSpPr', txBox='1')
                etree.SubElement(nvSpPr2, f'{{{p_ns}}}nvPr')

                spPr2 = etree.SubElement(content_sp, f'{{{p_ns}}}spPr')
                xfrm2 = etree.SubElement(spPr2, f'{{{a}}}xfrm')
                etree.SubElement(xfrm2, f'{{{a}}}off', x='781050', y='1786255')
                etree.SubElement(xfrm2, f'{{{a}}}ext', cx='10805160', cy='4500000')
                etree.SubElement(xfrm2, f'{{{a}}}prstGeom2', prst='rect')
                etree.SubElement(etree.SubElement(content_sp, f'{{{p_ns}}}txBody'), f'{{{a}}}bodyPr') # 占位
                txBody2 = content_sp.find(f'.//{{{p_ns}}}txBody') # 重定向

            # 遍历添加段落
            for item in page_items:
                txt = item.get('text', '')
                level = item.get('level', 0)
                if txt:
                    para = etree.SubElement(txBody2, f'{{{a}}}p')
                    pPr = etree.SubElement(para, f'{{{a}}}pPr')
                    if level >= 1:
                        # 2级及以上：开启 Bullet 项目符号 和 缩进
                        pPr.set('marL', '571500')
                        etree.SubElement(pPr, f'{{{a}}}buFont', typeface='Arial')
                        etree.SubElement(pPr, f'{{{a}}}buChar', char='•')
                    
                    lnSpc = etree.SubElement(pPr, f'{{{a}}}lnSpc')
                    etree.SubElement(lnSpc, f'{{{a}}}spcPct', val='150000')  # 1.5倍
                    run = etree.SubElement(para, f'{{{a}}}r')
                    rPr = etree.SubElement(run, f'{{{a}}}rPr', lang='zh-CN', sz=str(font_sz_val), dirty='0')
                    solidFill = etree.SubElement(rPr, f'{{{a}}}solidFill')
                    etree.SubElement(solidFill, f'{{{a}}}schemeClr', val='tx1')
                    etree.SubElement(rPr, f'{{{a}}}latin', typeface='Microsoft YaHei Light')
                    etree.SubElement(rPr, f'{{{a}}}ea', typeface='Microsoft YaHei Light')
                    t = etree.SubElement(run, f'{{{a}}}t')
                    t.text = txt

            # 追加专属导航条 (含 prs)
            add_nav_bar_to_slide(new_slide, code_str, all_chapters, prs)

    # 3. ---- B. 图片独占单页 ----
    from urllib.request import urlretrieve  # 确保可用
    def download_image_local(uri):
        if not uri: return None
        if uri.startswith('http'): url = uri
        else: url = f'https://mubu.com/{uri}'
        os.makedirs('/tmp/outline_images', exist_ok=True)
        local_filename = os.path.join('/tmp/outline_images', os.path.basename(uri))
        if os.path.exists(local_filename): return local_filename
        try: urlretrieve(url, local_filename); return local_filename
        except Exception: return None

    for img in all_images:
        image_uri = img.get('uri')
        if image_uri:
            local_img = download_image_local(image_uri)
            if local_img:
                try:
                    img_slide = prs.slides.add_slide(default_layout)
                    generated_slides.append(img_slide)

                    # 大图页标题同样置空
                    if hasattr(img_slide.shapes, 'title') and img_slide.shapes.title:
                        img_slide.shapes.title.text = page_title
                    else:
                        for shp in img_slide.shapes:
                            if hasattr(shp, 'has_text_frame') and shp.has_text_frame:
                                if '标题' in shp.text_frame.text or shp.name == '标题':
                                    shp.text_frame.text = page_title
                                    break
                    
                    slide_w, slide_h = prs.slide_width, prs.slide_height
                    top_offset = Emu(1786255)
                    margin = Emu(500000)
                    avail_h, avail_w = slide_h - top_offset - margin, slide_w - margin * 2
                    
                    fitted_h = avail_h
                    fitted_w = int(avail_h * (img.get('w', 1470) / (img.get('h', 786) or 1)))
                    if fitted_w > avail_w:
                        fitted_w = avail_w
                        fitted_h = int(avail_w * (img.get('h', 786) / (img.get('w', 1470) or 1)))
                    left = (slide_w - fitted_w) // 2
                    top = top_offset + (avail_h - fitted_h) // 2
                    img_slide.shapes.add_picture(local_img, left, top, width=fitted_w, height=fitted_h)

                    # 追加专属导航条 (含 prs)
                    add_nav_bar_to_slide(img_slide, code_str, all_chapters, prs)
                except Exception as e:
                    print(f"追加大图页失败: {e}")

    return generated_slides

    return generated_slides


# ==================== 主函数 ====================
def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    input_path = os.path.join(base_dir, 'mubu_response.json')
    template_path = os.path.join(base_dir, 'tpl', 'PPT模板.pptx')
    output_path = os.path.join(base_dir, 'output', 'presentation.pptx')

    # 确保输出目录存在
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    print("加载 v2版 幕布大纲树...")
    from mubu_parser import load_mubu_data_v2
    course_data = load_mubu_data_v2(input_path)
    
    if not course_data:
        print("错误: 无法加载幕布数据")
        return

    summary_cn = course_data.get('courseName', '')
    summary_en = course_data.get('chaptersEnName', '')
    outline = course_data # 整个 course_data 传给后续，或者直接传 chapters
    
    # 为了让 create_toc_slide 等函数能直接吃到 chapters
    chapters = course_data.get('chapters', [])
    print(f"解析大纲完成：{len(chapters)} 个大章")
    print(f"文档标题摘要(中): {summary_cn}")
    print(f"文档标题摘要(英): {summary_en}")

    # ---- 0. 中间生成 JSON 方便分析 ----
    debug_json_v2_path = os.path.join(base_dir, 'output', 'mubu_parsed_structure_v2_slides.json')
    try:
        with open(debug_json_v2_path, 'w', encoding='utf-8') as f:
            json.dump(course_data, f, ensure_ascii=False, indent=2)
        print(f"新版中间解析数据已导出至: {debug_json_v2_path}")
    except Exception as e:
        print(f"导出中间 JSON 失败: {e}")
        print(f"导出中间 JSON 失败: {e}")

    # 打开已生成母版的模板
    prs = Presentation(template_path)



    # 保存引用
    slides = list(prs.slides)
    cover_slide = slides[0]   # 封面

    # ---- 1. 替换封面文字 ----
    print("替换封面页文字...")
    import math
    for shp in cover_slide.shapes:
        if shp.has_text_frame:
            for paragraph in shp.text_frame.paragraphs:
                for run in paragraph.runs:
                    txt = run.text.strip()
                    # A. 严格全字匹配，防止“课程名称英文”包含“课程名称”而造成误伤
                    if txt == '课程名称' or txt == 'XXXXX（课程名称）':
                        run.text = summary_cn
                    elif txt == '课程名称英文' or txt == 'XXXXXXXXXXX（课程名称英文）':
                        run.text = summary_en
                    elif '课程名称' in txt or 'XXXXX' in txt:
                        # B. 模糊替换
                        new_txt = run.text.replace('XXXXX（课程名称）', summary_cn) \
                                         .replace('XXXXXXXXXXX（课程名称英文）', summary_en) \
                                         .replace('课程名称英文', summary_en) \
                                         .replace('课程名称', summary_cn) \
                                         .replace('XXXXX', '') \
                                         .replace('XXXXXXXXXXX', '')
                        run.text = new_txt

    # 封面的“免责声明”文本框防重叠自适应位移
    # 估算主标题行数 (中文汉字约14字一行, 英文30字符一行)
    cn_lines = math.ceil(len(summary_cn) / 14)
    en_lines = math.ceil(len(summary_en) / 30)
    extra = max(0, cn_lines - 1) + max(0, en_lines - 1)

    if extra > 0:
        offset = extra * 420000 # EMU单位位移值
        for shp in cover_slide.shapes:
            if shp.has_text_frame and '未经授权严禁' in shp.text_frame.text:
                shp.top = shp.top + offset
                print(f"  [封面自适应] 检测到多出 {extra} 行文字，下方框向下平移 {offset} EMU")

    toc_slide = slides[1]         # 总目录
    chapter_template = slides[2]  # 章节页模板
    content_template = slides[3]  # 内容页模板
    ending_slide = slides[4]  # 结束页

    # 记住布局
    chapter_layout = chapter_template.slide_layout  # 1_标题和内容 (母版0)
    content_layout = content_template.slide_layout   # 标题和竖排文字 (母版1)

    print("修改总目录页...")
    create_toc_slide(prs, toc_slide, chapters)

    # 找到母版对应的布局
    # 新母版从索引5开始（0-4是原始5个母版）
    # 母版1-1-1对应slideMaster6，其布局是slideLayout23
    # 我们需要找到这些布局
    # 由于新母版的布局是我们创建的，每个母版只有一个布局
    # 我们需要通过布局名称来找到对应的布局

    # 收集所有可用的布局（按母版名称索引）
    master_layouts = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            cSld = layout._element.find(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
            if cSld is not None:
                name = cSld.get('name', '')
                if name.startswith('母版'):
                    master_layouts[name] = layout

    print(f"找到 {len(master_layouts)} 个母版布局")

    # 现在生成所有内容页
    # 先删除原始的第3、4页（章节页和内容页模板），保留封面、总目录和结束页
    # 然后按顺序添加：章节页 + 内容页

    # 现在生成所有内容页
    # 递归展开 Chapter 树，适配无限级联导航
    def expand_chapter_slides(prs, chapter, level_path):
        slides_added = []
        ch_name = chapter.get('chapterName', '')
        level = chapter.get('level', 1)

        # A. 如果是 Level 1 (章节)，添加章节页
        if level == 1:
            ch_idx = level_path[0]
            print(f"  添加章节页: {ch_name}")
            ch_slide = create_chapter_slide(prs, chapter_template, chapter, ch_idx)
            slides_added.append(('chapter', ch_slide))

        # B. 如果包含 content 章节内容，添加内容页
        contents = chapter.get('content', [])
        if contents:
            # 拼合母版定位名, 比如 "1.1.2" -> "母版1-1-2"
            code_str = chapter.get('code', '1.1.1')
            parts = code_str.split('.')
            while len(parts) < 3:
                parts.append('1') # 补齐三级供旧模板向下兼容
            master_name = f'母版{"-".join(parts[:3])}'
            
            layout = master_layouts.get(master_name, content_layout)
            print(f"  添加内容页: {ch_name} -> {master_name} (包含 {len(contents)} 项内容)")

            # 调用新版内容页生成器
            code_str_raw = chapter.get('code', '1.1.1')
            content_slides = create_content_slide_v2(prs, layout, ch_name, contents, content_layout, code_str_raw, top_chapters)
            for cs in content_slides:
                slides_added.append(('content', cs))

        # C. 递归遍历子章节
        for i, sub_ch in enumerate(chapter.get('subChapter', [])):
             sub_slides = expand_chapter_slides(prs, sub_ch, level_path + [i])
             slides_added.extend(sub_slides)

        return slides_added

    # 循环顶级 chapters 执行递归投喂
    new_slides = []
    # 从 outline 中提取顶级章
    top_chapters = outline.get('chapters', []) if isinstance(outline, dict) else outline
    
    for ch_idx, chapter in enumerate(top_chapters):
        new_slides.extend(expand_chapter_slides(prs, chapter, [ch_idx]))

    # 重新排序幻灯片
    # 目标顺序：封面(0), 总目录(1), [新页...], 结束页
    # 当前顺序：封面(0), 总目录(1), 章节模板(2), 内容模板(3), 结束页(4), [新页...]

    # 通过 XML 操作重排幻灯片顺序
    prs_element = prs.element if hasattr(prs, 'element') else prs.part.element
    
    sldIdLst = prs_element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')

    # 获取所有 sldId 元素
    all_sldIds = list(sldIdLst)

    # 原始5页 + 新增页
    # 我们要删除索引2和3（章节模板和内容模板），然后把新页放在结束页之前
    # all_sldIds[0] = 封面
    # all_sldIds[1] = 总目录
    # all_sldIds[2] = 章节模板 -> 删除
    # all_sldIds[3] = 内容模板 -> 删除
    # all_sldIds[4] = 结束页
    # all_sldIds[5:] = 新页

    cover_id = all_sldIds[0]
    toc_id = all_sldIds[1]
    ending_id = all_sldIds[4]
    new_page_ids = all_sldIds[5:]  # 新添加的页
    old_template_ids = [all_sldIds[2], all_sldIds[3]]  # 要删除的模板页

    # 清空 sldIdLst
    for child in list(sldIdLst):
        sldIdLst.remove(child)

    # 重新按顺序添加
    sldIdLst.append(cover_id)
    sldIdLst.append(toc_id)
    for sid in new_page_ids:
        sldIdLst.append(sid)
    sldIdLst.append(ending_id)

    # 注意：我们不删除模板页的实际slide数据，只是不在sldIdLst中引用
    # 这样它们在打开时不会显示

    total_pages = 2 + len(new_page_ids) + 1  # 封面 + 总目录 + 新页 + 结束
    print(f"\n生成完成！共 {total_pages} 页")
    print(f"  封面: 1页")
    print(f"  总目录: 1页")
    print(f"  章节页: {sum(1 for t, _ in new_slides if t == 'chapter')} 页")
    print(f"  内容页: {sum(1 for t, _ in new_slides if t == 'content')} 页")
    print(f"  结束页: 1页")

    prs.save(output_path)
    print(f"\n保存成功: {output_path}")


if __name__ == '__main__':
    main()

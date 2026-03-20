#!/usr/bin/env python3
import os
from pptx import Presentation

tpl_path = "tpl/PPT模板.pptx"
final_path = "output/final_presentation.pptx"

print("=== 1. 原模板第4页 文字样式 ===")
if os.path.exists(tpl_path):
    prs_tpl = Presentation(tpl_path)
    slide4 = prs_tpl.slides[3]
    for i, shp in enumerate(slide4.shapes):
        if shp.has_text_frame:
            print(f"Shape {i} [{shp.name}]:")
            for p_idx, p in enumerate(shp.text_frame.paragraphs):
                 for r_idx, r in enumerate(p.runs):
                      print(f"  P[{p_idx}] R[{r_idx}] '{r.text}'")
                      if r.font:
                           sz = r.font.size.pt if r.font.size else "None"
                           print(f"    Font: {r.font.name}, Size: {sz}pt, Bold: {r.font.bold}")

print("\n=== 2. 最终 PPT 特定页母版检视 ===")
if os.path.exists(final_path):
    prs_final = Presentation(final_path)
    target_indices = [16, 17, 18, 19, 32, 33, 48, 49] # 0-based indices for 17-20, 33-34, 49-50
    for idx in target_indices:
        if idx < len(prs_final.slides):
             slide = prs_final.slides[idx]
             lay = slide.slide_layout
             print(f"\n--- Slide {idx+1} (Layout: {lay.name}) ---")
             for shp in slide.shapes:
                  if shp.has_text_frame:
                       txt = shp.text_frame.text.replace('\n', ' ').strip()
                       print(f"  Shape: {txt[:60]}...")
        else:
             print(f"Slide {idx+1} out of bounds.")

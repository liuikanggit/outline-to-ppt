#!/usr/bin/env python3
from pptx import Presentation

pptx_path = "output/master_with_toc.pptx"
prs = Presentation(pptx_path)

if len(prs.slides) >= 4:
    slide = prs.slides[3] # 第4页
    print("=== Slide 4 Placeholders ===")
    for ph in slide.placeholders:
        print(f"Type: {ph.placeholder_format.type}")
        print(f"  Name: {ph.name}")
        print(f"  Left: {ph.left}")
        print(f"  Top: {ph.top}")
        print(f"  Width: {ph.width}")
        print(f"  Height: {ph.height}")
else:
    print("Slide length too short")

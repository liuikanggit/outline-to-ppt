#!/usr/bin/env python3
import os
from pptx import Presentation

pptx_path = "output/master_with_toc.pptx"
prs = Presentation(pptx_path)

# 打印第一个加进去的母版（应该是母版-1.1 或者类似）
for m in prs.slide_masters:
    for lay in m.slide_layouts:
        if lay.name.startswith("母版") or "SlideLayout" in lay.name:
            print(f"\n=== Layout Name: {lay.name} ===")
            print(f"Total shapes in layout: {len(lay.shapes)}")
            for i, shp in enumerate(lay.shapes):
                print(f"  Shape {i}: Type {type(shp)}")
                if shp.has_text_frame:
                    print(f"    Text: '{shp.text_frame.text}'")

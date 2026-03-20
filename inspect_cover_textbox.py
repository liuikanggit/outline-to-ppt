#!/usr/bin/env python3
import os
from pptx import Presentation

def inspect_disclaimer(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return
        
    print(f"\n--- {os.path.basename(pptx_path)} ---")
    try:
        prs = Presentation(pptx_path)
        slide = prs.slides[0]
        for i, shp in enumerate(slide.shapes):
            if shp.has_text_frame:
                txt = shp.text_frame.text
                if '未经授权严禁' in txt:
                    print(f"Found Disclaimer TextBox:")
                    print(f"  Shape Index: {i}")
                    print(f"  Shape Name: {shp.name}")
                    print(f"  Top: {shp.top} EMU")
                    print(f"  Left: {shp.left} EMU")
                    print(f"  Width: {shp.width} EMU")
                    print(f"  Height: {shp.height} EMU")
                    # 顺便看下文本框内的内容和段落数
                    print(f"  Text: {txt[:50]}...")
            # 顺便打印下主标题的 TextBox 位置和高度，方便对比遮挡
            if shp.has_text_frame and ('伟经集团' in shp.text_frame.text or '课程名称' in shp.text_frame.text):
                print(f"Found Title TextBox:")
                print(f"  Shape Index: {i}")
                print(f"  Top: {shp.top} EMU")
                print(f"  Height: {shp.height} EMU")
                print(f"  Text: {shp.text_frame.text[:50]}...")
    except Exception as e:
        print(f"Error inspecting {pptx_path}: {e}")

inspect_disclaimer("tpl/PPT模板.pptx")
inspect_disclaimer("output/master.pptx")
inspect_disclaimer("output/master_with_cover.pptx")

#!/usr/bin/env python3
from pptx import Presentation

pptx_path = "tpl/PPT模板.pptx"
prs = Presentation(pptx_path)

if len(prs.slides) >= 4:
    slide = prs.slides[3]  # 第4页 (0-indexed 3)
    print("=== Original Slide 4 Standard Shapes ===")
    for i, shp in enumerate(slide.shapes):
        print(f"Shape {i}: {shp.name}")
        if shp.has_text_frame:
            print(f"  Text: '{shp.text_frame.text}'")
            print(f"  Left: {shp.left} EMU")
            print(f"  Top: {shp.top} EMU")
            print(f"  Width: {shp.width} EMU")
            print(f"  Height: {shp.height} EMU")
else:
    print("Template too short")

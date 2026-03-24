#!/usr/bin/env python3
"""
第6步骤：生成 PPT 正文页（章节页 + 内容页）
根据 mubu_parsed_structure.json 的大纲层次结构，递归加入幻灯片。
内容页使用事先做好的、带联动导航条的 SlideMaster / Layout。
"""

import json
import os
import math
from pathlib import Path
from urllib.request import urlretrieve
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

# ----------------- Open XML 命名空间 -----------------
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

# ----------------- 预设第4页 文本框 物理坐标 (从原模版第4页提取) -----------------
TITLE_L = 664845
TITLE_T = 1242060
TITLE_W = 6096000
TITLE_H = 368300

BODY_L = 664845
BODY_T = 1708785
BODY_W = 10805160
BODY_H = 4546600

def download_image_local(uri):
    """下载图片缓存到本地"""
    if not uri: return None
    if uri.startswith('http'): url = uri
    else: url = f'https://mubu.com/{uri}'
    temp_dir = Path('/tmp/outline_images')
    temp_dir.mkdir(parents=True, exist_ok=True)
    local_filename = temp_dir / os.path.basename(uri)
    if local_filename.exists(): return str(local_filename)
    try: 
        urlretrieve(url, str(local_filename))
        return str(local_filename)
    except Exception: return None

# ----------------- 1. 自动分页与字号算法 -----------------
def paginate_contents_custom(contents):
    """
    符合用户直觉的智能防打散分页算法：
    将列表转化为节点树，每个节点计算其子树所需行数(block_lines)。
    1. 若某节点及其所有子节点能放在同一页，绝不拆断。
    2. 若当前页放不下，但放入新页能完整放下，则提前开新页。
    3. 只有当节点族群超出 14 行时，才会拆散树结构强制降维分布。
    """
    def get_lines(txt, max_char_per_line=59):
        return math.ceil(len(txt) / max_char_per_line) if txt else 1

    def build_tree(item, depth):
        txt = item.get('text', '')
        has_bullet = item.get('hasBullet', False)
        node = {
            'text': txt,
            'level': depth,
            'hasBullet': has_bullet or (depth >= 1),
            'lines': get_lines(txt, 59), 
            'children': []
        }
        for sub in item.get('subContent', []):
            node['children'].append(build_tree(sub, depth + 1))
        
        node['block_lines'] = node['lines'] + sum(c['block_lines'] for c in node['children'])
        return node
    
    forest = []
    for c in contents:
        if c.get('type') == 'text':
            info = c.get('text', {})
            info_list = info if isinstance(info, list) else [info]
            for node_data in info_list:
                if node_data.get('text', ''):
                    forest.append(build_tree(node_data, 0))

    if not forest:
        return [], 1600

    total_lines = sum(n['block_lines'] for n in forest)
    
    # 将树形直接摊平的辅助函数
    def flatten_tree(nodes):
        flat = []
        def dfs(n):
            flat.append(n)
            for ch in n['children']: dfs(ch)
        for n in nodes: dfs(n)
        return flat

    if total_lines <= 12:
        return [flatten_tree(forest)], 1600
    if total_lines <= 14:
        return [flatten_tree(forest)], 1400

    # 突破限制，进行智能组装分页
    pages = []
    curr_page = []
    curr_lines = 0
    MAX_LINES = 14

    def pack_node(node):
        nonlocal pages, curr_page, curr_lines
        
        # 1. 尝试将整个节点群（自身 + 所有子孙）整体放入当前页
        if curr_lines + node['block_lines'] <= MAX_LINES:
            def dfs_add(n):
                nonlocal curr_lines
                curr_page.append(n)
                curr_lines += n['lines']
                for ch in n['children']: dfs_add(ch)
            dfs_add(node)
            return

        # 2. 装不下当前页，判断放到新页能否完整装下
        if node['block_lines'] <= MAX_LINES:
            if curr_page:
                pages.append(curr_page)
            curr_page = []
            curr_lines = 0
            def dfs_add2(n):
                nonlocal curr_lines
                curr_page.append(n)
                curr_lines += n['lines']
                for ch in n['children']: dfs_add2(ch)
            dfs_add2(node)
            return

        # 3. 连一整页都装不下这个节点群，必须强制拆分！先装自身
        if curr_lines + node['lines'] > MAX_LINES:
            if curr_page:
                pages.append(curr_page)
            curr_page = []
            curr_lines = 0
        
        curr_page.append(node)
        curr_lines += node['lines']

        # 然后对其子节点依次调用包装逻辑
        for ch in node['children']:
            pack_node(ch)

    for n in forest:
        pack_node(n)
        
    if curr_page:
        pages.append(curr_page)

    return pages, 1600


# ----------------- 2. 章节页生成 -----------------
def create_chapter_slide(prs, template_slide, chapter):
    """
    复制模板页（第3页），修改文本内容
    """
    layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    
    # 彻底清除由 layout 带来的其它非 placeholders 形状
    for ph in list(new_slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    spTree_new = new_slide.shapes._spTree
    spTree_old = template_slide.shapes._spTree
    
    for sp in spTree_old:
        tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
        if tag == 'sp':
            new_sp = etree.fromstring(etree.tostring(sp))
            spTree_new.append(new_sp)
            
            # 修改内容
            txBody = new_sp.find(f'.//{{{NSMAP["p"]}}}txBody')
            if txBody is not None:
                paras = txBody.findall(f'.//{{{NSMAP["a"]}}}p')
                if len(paras) >= 1:
                    # 清空段落下所有 a:r
                    r_nodes = paras[0].findall(f'.//{{{NSMAP["a"]}}}r')
                    if r_nodes:
                         t_node = r_nodes[0].find(f'{{{NSMAP["a"]}}}t')
                         if t_node is not None: t_node.text = chapter.get('chapterName', '')
                         # 移除多余的 r 节点(例如残留的“名称”)
                         for extra_r in r_nodes[1:]:
                              paras[0].remove(extra_r)
                if len(paras) >= 2:
                    r_nodes = paras[1].findall(f'.//{{{NSMAP["a"]}}}r')
                    if r_nodes:
                         r_node = r_nodes[0]
                         t_node = r_node.find(f'{{{NSMAP["a"]}}}t')
                         if t_node is not None: 
                              t_node.text = chapter.get('chapterEnName', '')
                         
                         # 将英文字体颜色覆写为灰色、字号设定为 18 号 (18*100 = 1800)
                         rPr = r_node.find(f'{{{NSMAP["a"]}}}rPr')
                         if rPr is None:
                             rPr = etree.Element(f'{{{NSMAP["a"]}}}rPr')
                             r_node.insert(0, rPr)
                         rPr.set('sz', '1800')
                         
                         # 清理旧的颜色节点
                         for child in list(rPr):
                             if child.tag.endswith('solidFill') or child.tag.endswith('gradFill'):
                                 rPr.remove(child)
                         
                         # 创建一个全新的纯灰色节点
                         solidFill = etree.Element(f'{{{NSMAP["a"]}}}solidFill')
                         etree.SubElement(solidFill, f'{{{NSMAP["a"]}}}srgbClr', val='808080')
                         
                         # 必须将其插在开头位置 (严格的 OpenXML Schema 规范: Fill 系必须排在 latin/ea 之内)
                         insert_idx = 0
                         for idx, child in enumerate(rPr):
                             if child.tag.endswith('ln'):
                                 insert_idx = idx + 1
                         rPr.insert(insert_idx, solidFill)

                         for extra_r in r_nodes[1:]:
                              paras[1].remove(extra_r)

    return new_slide


# ----------------- 3. 内容页生成 -----------------
def create_content_slide_optimized(prs, content_layout, page_title, items, font_sz_val):
    """
    创建一个内容页，填充文字（基于固定预设坐标，规避 Blank Layout 无 Placeholder 问题）
    """
    new_slide = prs.slides.add_slide(content_layout)
    
    # 1. 创建标题文本框并填充
    title_box = new_slide.shapes.add_textbox(TITLE_L, TITLE_T, TITLE_W, TITLE_H)
    t_frame = title_box.text_frame
    p = t_frame.paragraphs[0]
    r = p.add_run()
    r.text = page_title
    # 为 Title 应用粗体及微软雅黑字体族风格
    rPr = r._r.get_or_add_rPr()
    rPr.set('b', '1') # Bold
    etree.SubElement(rPr, f'{{{NSMAP["a"]}}}latin', typeface='微软雅黑')
    etree.SubElement(rPr, f'{{{NSMAP["a"]}}}ea', typeface='微软雅黑')
    
    # 2. 创建正文文本框
    body_box = new_slide.shapes.add_textbox(BODY_L, BODY_T, BODY_W, BODY_H)
    body_box.text_frame.word_wrap = True
    
    txBody = body_box._element.find('.//' + f'{{{NSMAP["p"]}}}txBody')
    if txBody is not None:
         # 清空 python-pptx 默认创建的空段落，防止首行空行
         for p in txBody.findall(f'.//{{{NSMAP["a"]}}}p'):
              txBody.remove(p)
    else:
         txBody = etree.SubElement(body_box._element, f'{{{NSMAP["p"]}}}txBody')
         etree.SubElement(txBody, f'{{{NSMAP["a"]}}}bodyPr')
         etree.SubElement(txBody, f'{{{NSMAP["a"]}}}lstStyle')

    # 填充 paragraphs
    for item in items:
        # 兼容 string 和 dict
        txt = item.get('text', '')
        level = item.get('level', 0)
        has_bullet = item.get('hasBullet', False)

        if txt:
            para = etree.SubElement(txBody, f'{{{NSMAP["a"]}}}p')
            pPr = etree.SubElement(para, f'{{{NSMAP["a"]}}}pPr')
            
            # PowerPoint 多级缩进字典 (marL, indent, char) — 统一使用实心圆点
            bullet_styles = [
                (342900, -342900, '•'),  # L0
                (742950, -342900, '•'),  # L1
                (1143000, -342900, '•'), # L2
                (1543050, -342900, '•'), # L3
                (1943100, -342900, '•')  # L4
            ]
            
            # 根据扁平化下传的 depth 获取专属缩进配置 (兼容越界场景)
            lvl_idx = min(level, len(bullet_styles) - 1)

            # ==========================
            # 严格的 OpenXML Schema 顺序:
            # lnSpc -> ... -> buFont -> buChar
            # ==========================
            lnSpc = etree.SubElement(pPr, f'{{{NSMAP["a"]}}}lnSpc')
            etree.SubElement(lnSpc, f'{{{NSMAP["a"]}}}spcPct', val='150000')  # 1.5倍行距

            if has_bullet:
                marL, ind_val, char = bullet_styles[lvl_idx]
                pPr.set('marL', str(marL))
                pPr.set('indent', str(ind_val))
                etree.SubElement(pPr, f'{{{NSMAP["a"]}}}buFont', typeface='Arial')
                etree.SubElement(pPr, f'{{{NSMAP["a"]}}}buChar', char=char)
            elif level >= 1:
                # 仅缩进，无项目符号
                marL, _, _ = bullet_styles[lvl_idx]
                pPr.set('marL', str(marL))

            run = etree.SubElement(para, f'{{{NSMAP["a"]}}}r')
            rPr = etree.SubElement(run, f'{{{NSMAP["a"]}}}rPr', lang='zh-CN', sz=str(font_sz_val), dirty='0')
            solidFill = etree.SubElement(rPr, f'{{{NSMAP["a"]}}}solidFill')
            etree.SubElement(solidFill, f'{{{NSMAP["a"]}}}schemeClr', val='tx1')
            etree.SubElement(rPr, f'{{{NSMAP["a"]}}}latin', typeface='Microsoft YaHei Light')
            etree.SubElement(rPr, f'{{{NSMAP["a"]}}}ea', typeface='Microsoft YaHei Light')
            t = etree.SubElement(run, f'{{{NSMAP["a"]}}}t')
            t.text = txt

    return new_slide

def append_images_to_slide(new_slide, img_info, prs):
    """独占单页插入大图：高度填满且自适应"""
    image_uri = img_info.get('uri')
    local_img = download_image_local(image_uri)
    if not local_img: 
         return None

    try:
        img_w_orig = img_info.get('w', 1470)
        img_h_orig = img_info.get('h', 786) or 1

        # 高度填满，宽度自适应 (基于预设正文范围)
        fitted_h = BODY_H
        fitted_w = int(BODY_H * (img_w_orig / img_h_orig))

        if fitted_w > BODY_W:
             fitted_w = BODY_W
             fitted_h = int(BODY_W * (img_h_orig / img_w_orig))

        left = BODY_L + (BODY_W - fitted_w) // 2
        top = BODY_T + (BODY_H - fitted_h) // 2
        
        new_slide.shapes.add_picture(local_img, left, top, width=fitted_w, height=fitted_h)
        return True
    except Exception as e:
         print(f"  [图片报错] {e}")
         return False


# ----------------- 4. 路由逻辑 -----------------
def get_layout_by_names(prs, target_master_name, target_layout_name, default_layout):
    """
    匹配重构后的多版面逻辑：通过 masterName + layoutName 双键联合查找对应的 Layout
    """
    # 第一次精确匹配：masterName 且 layoutName 相同
    for m in prs.slide_masters:
        m_cSld = m._element.find(f'{{{NSMAP["p"]}}}cSld')
        m_name = m_cSld.get('name', '') if m_cSld is not None else ''
        
        if m_name == target_master_name:
            for lay in m.slide_layouts:
                lay_cSld = lay._element.find(f'{{{NSMAP["p"]}}}cSld')
                lay_name = lay_cSld.get('name', '') if lay_cSld is not None else ''
                if lay_name == target_layout_name:
                    return lay
                    
    # 第二次降级匹配：如果没找到特定 layoutName，则降级查找相同 masterName 下的 'default'
    for m in prs.slide_masters:
        m_cSld = m._element.find(f'{{{NSMAP["p"]}}}cSld')
        m_name = m_cSld.get('name', '') if m_cSld is not None else ''
        if m_name == target_master_name:
            for lay in m.slide_layouts:
                lay_cSld = lay._element.find(f'{{{NSMAP["p"]}}}cSld')
                lay_name = lay_cSld.get('name', '') if lay_cSld is not None else ''
                if lay_name == "default":
                    return lay
                    
    # 都没找到，返回骨架 fallback
    return default_layout

def main():
    import argparse
    parser = argparse.ArgumentParser(description="根据大纲 JSON 递归生成正文页")
    parser.add_argument("--input_json", default="output/mubu_parsed_structure.json", help="大纲 JSON 路径")
    parser.add_argument("--input_pptx", default="output/master_with_toc.pptx", help="第5部输出已带总目录的 PPTX 路径")
    parser.add_argument("--output_pptx", default="output/final_presentation.pptx", help="最终合流成品的命名路径")
    args = parser.parse_args()

    input_json = Path(args.input_json)
    input_pptx = Path(args.input_pptx)
    output_pptx = Path(args.output_pptx)

    if not input_json.exists() or not input_pptx.exists():
         print("❌ 缺少配套 input_json 或 input_pptx 文件。")
         return

    with open(input_json, "r", encoding="utf-8") as f:
         data = json.load(f)
    outline_chapters = data.get("chapters", [])

    prs = Presentation(input_pptx)
    slides = list(prs.slides)
    if len(slides) < 4:
         print("❌ 输入 PPT 的内容不包含章节页/内容页骨架，无法向下级联。")
         return

    chapter_template = slides[2]  # 章节页骨架 template
    content_template = slides[3]  # 内容页默认 fallback
    default_content_layout = content_template.slide_layout

    added_slides = []

    def expand_chapters(chapter, level_path, ancestor_names):
        ch_name = chapter.get('chapterName', '')
        level = chapter.get('level', 1)
        code_str = chapter.get('code', '1.1.1')
        
        current_names = ancestor_names + [ch_name]
        
        target_master_name = ""
        target_layout_name = ""
        # 兼容最新母版分组生成规则
        if level == 1:
            target_master_name = current_names[0]
            target_layout_name = "default"
        elif level == 2:
            target_master_name = f"{current_names[0]}-{current_names[1]}"
            target_layout_name = "default"
        elif level >= 3:
            target_master_name = f"{current_names[0]}-{current_names[1]}"
            target_layout_name = current_names[2]


        # 1. 章节页
        if level == 1:
            print(f"🌲 生成[章节页]: {ch_name}")
            ch_slide = create_chapter_slide(prs, chapter_template, chapter)
            added_slides.append(ch_slide)

        # 2. 内容页 (针对包含 content)
        contents = chapter.get('content', [])
        if contents:
            print(f"📄 生成[内容页]: {ch_name} (包含 {len(contents)} 个节点)")
            text_contents = [c for c in contents if c.get('type') == 'text']
            img_contents = [c for c in contents if c.get('type') == 'image']

            if text_contents:
                # 恢复全量内容页渲染
                pages, font_sz = paginate_contents_custom(text_contents)
                lay = get_layout_by_names(prs, target_master_name, target_layout_name, default_content_layout)
                for page_dict in pages:
                    pg_slide = create_content_slide_optimized(prs, lay, ch_name, page_dict, font_sz)
                    added_slides.append(pg_slide)

            for img in img_contents:
                lay = get_layout_by_names(prs, target_master_name, target_layout_name, default_content_layout)
                img_slide = prs.slides.add_slide(lay)
                
                # 独家修复：图片页强制增加独立标题框
                title_box = img_slide.shapes.add_textbox(TITLE_L, TITLE_T, TITLE_W, TITLE_H)
                p = title_box.text_frame.paragraphs[0]
                r = p.add_run()
                r.text = ch_name
                rPr = r._r.get_or_add_rPr()
                rPr.set('b', '1')
                etree.SubElement(rPr, f'{{{NSMAP["a"]}}}latin', typeface='微软雅黑')
                etree.SubElement(rPr, f'{{{NSMAP["a"]}}}ea', typeface='微软雅黑')
                
                append_images_to_slide(img_slide, img.get('image', {}), prs)
                added_slides.append(img_slide)

        for sub in chapter.get('subChapter', []):
             expand_chapters(sub, level_path + [sub.get('code', '')], current_names)

    print("✍️ 开始递归解析和挂载正文页...")
    for ch in outline_chapters:
         expand_chapters(ch, [], [])

    # 4. 抹杀多余模板页 [2、3页]
    prs_element = prs.element if hasattr(prs, 'element') else prs.part.element
    sldIdLst = prs_element.find(f'{{{NSMAP["p"]}}}sldIdLst')
    if sldIdLst is not None:
         all_sldIds = list(sldIdLst)
         if len(all_sldIds) > 4:
              cover_id = all_sldIds[0]
              toc_id = all_sldIds[1]
              ending_id = all_sldIds[4] if len(all_sldIds) > 4 else None
              # 收集新增页
              current_ids = [cover_id, toc_id]
              for cs in added_slides:
                   # 找到对应的新增 sldId 引用不太安全
                   pass
              # 还是按照索引重新构造安全
              new_ids = all_sldIds[5:]
              for aid in list(sldIdLst):
                   sldIdLst.remove(aid)
              sldIdLst.append(cover_id)
              sldIdLst.append(toc_id)
              for nid in new_ids:
                   sldIdLst.append(nid)
              if ending_id is not None:
                   sldIdLst.append(ending_id)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    print(f"✅ 正文页级联生成完成 -> {output_pptx}")

if __name__ == "__main__":
    main()

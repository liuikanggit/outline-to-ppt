#!/usr/bin/env python3
import os
from pptx import Presentation

pptx_path = "output/master_with_toc.pptx"

if not os.path.exists(pptx_path):
    print("File not found.")
    exit()

prs = Presentation(pptx_path)
count = 0
for m in prs.slide_masters:
    for lay in m.slide_layouts:
        for shp in lay.shapes:
            if shp.has_text_frame:
                txt = shp.text_frame.text
                if txt in ["标题", "文字内容"]:
                    print(f"Found on layout [{lay.name}]: {txt}")
                    shp.text_frame.text = "" # 尝试擦除
                    count += 1

print(f"Total placeholders cleared: {count}")
prs.save("output/master_cleared.pptx")
print("Saved testing cleared file.")

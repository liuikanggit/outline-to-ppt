#!/usr/bin/env python3
import os
from pptx import Presentation

tpl_path = "tpl/PPT模板.pptx"

print("=== 1. 原模板第4页 文字样式 ===")
if os.path.exists(tpl_path):
    prs_tpl = Presentation(tpl_path)
    slide4 = prs_tpl.slides[3]
    for i, shp in enumerate(slide4.shapes):
        if shp.has_text_frame:
            print(f"\nShape {i} [{shp.name}]:")
            for p_idx, p in enumerate(shp.text_frame.paragraphs):
                 for r_idx, r in enumerate(p.runs):
                      print(f"  P[{p_idx}] R[{r_idx}] '{r.text}'")
                      if r.font:
                           sz = r.font.size.pt if r.font.size else "None"
                           # 检查字体族
                           name = r.font.name
                           bold = r.font.bold
                           print(f"    Font: {name}, Size: {sz}pt, Bold: {bold}")
else:
    print("Template not found.")

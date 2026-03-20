#!/usr/bin/env python3
"""
第5步骤：自动生成 PPT 总目录页（第2页）
根据 mubu_parsed_structure.json 的 chapters（一级章节）横向排布生成总目录。
在第4步骤(generate_cover.py)输出的 PPTX 基础上执行。
"""

import json
import os
from pathlib import Path
from lxml import etree
from pptx import Presentation

# ----------------- Open XML 命名空间 -----------------
PML_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

def create_toc_chapter_box(chapter_dict, ch_idx, x, y, w, h, sp_id):
    """
    创建总目录中的章节标题文本框 (使用 etree)
    """
    sp = etree.Element(f'{{{PML_NS}}}sp')
    
    # 1. nvSpPr (形状信息)
    nvSpPr = etree.SubElement(sp, f'{{{PML_NS}}}nvSpPr')
    cNvPr = etree.SubElement(nvSpPr, f'{{{PML_NS}}}cNvPr', id=str(sp_id), name=f'TOC_Chapter_{ch_idx}')
    cNvSpPr = etree.SubElement(nvSpPr, f'{{{PML_NS}}}cNvSpPr', txBox='1')
    etree.SubElement(nvSpPr, f'{{{PML_NS}}}nvPr')

    # 2. spPr (样式与坐标)
    spPr = etree.SubElement(sp, f'{{{PML_NS}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{A_NS}}}xfrm')
    etree.SubElement(xfrm, f'{{{A_NS}}}off', x=str(x), y=str(y))
    etree.SubElement(xfrm, f'{{{A_NS}}}ext', cx=str(w), cy=str(h))
    prstGeom = etree.SubElement(spPr, f'{{{A_NS}}}prstGeom', prst='rect')
    etree.SubElement(prstGeom, f'{{{A_NS}}}avLst')

    # 3. txBody (文本内容)
    txBody = etree.SubElement(sp, f'{{{PML_NS}}}txBody')
    etree.SubElement(txBody, f'{{{A_NS}}}bodyPr', wrap='none', rtlCol='0')
    etree.SubElement(txBody, f'{{{A_NS}}}lstStyle')

    # 第一行：中文章节名（16pt 粗体 微软雅黑）
    para1 = etree.SubElement(txBody, f'{{{A_NS}}}p')
    etree.SubElement(para1, f'{{{A_NS}}}pPr', algn='ctr')
    run1 = etree.SubElement(para1, f'{{{A_NS}}}r')
    rPr1 = etree.SubElement(run1, f'{{{A_NS}}}rPr', lang='zh-CN', sz='1600', b='1', dirty='0')
    solidFill1 = etree.SubElement(rPr1, f'{{{A_NS}}}solidFill')
    etree.SubElement(solidFill1, f'{{{A_NS}}}schemeClr', val='tx1')
    etree.SubElement(rPr1, f'{{{A_NS}}}latin', typeface='微软雅黑')
    etree.SubElement(rPr1, f'{{{A_NS}}}ea', typeface='微软雅黑')
    run1_t = etree.SubElement(run1, f'{{{A_NS}}}t')
    run1_t.text = chapter_dict.get('chapterName', '')

    # 第二行：英文章节名（8pt Microsoft YaHei Light）
    name_en = chapter_dict.get('chapterEnName', '')
    if name_en:
        para2 = etree.SubElement(txBody, f'{{{A_NS}}}p')
        etree.SubElement(para2, f'{{{A_NS}}}pPr', algn='ctr')
        run2 = etree.SubElement(para2, f'{{{A_NS}}}r')
        rPr2 = etree.SubElement(run2, f'{{{A_NS}}}rPr', lang='en-US', sz='800', b='0', dirty='0')
        solidFill2 = etree.SubElement(rPr2, f'{{{A_NS}}}solidFill')
        etree.SubElement(solidFill2, f'{{{A_NS}}}schemeClr', val='tx1')
        etree.SubElement(rPr2, f'{{{A_NS}}}latin', typeface='Microsoft YaHei Light')
        etree.SubElement(rPr2, f'{{{A_NS}}}ea', typeface='Microsoft YaHei Light')
        run2_t = etree.SubElement(run2, f'{{{A_NS}}}t')
        run2_t.text = name_en

    return sp

def create_toc_slide(prs, slide, chapters):
    """
    修改页面，横向填充一级章节
    """
    spTree = slide.shapes._spTree
    shapes_to_remove = []

    # 1. 扫描出不需要的历史形状（除了叫“总目录”或叫“文本框 7”的框，其余全删）
    for sp in list(spTree):
        tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
        if tag == 'sp':
            cNvPr = sp.find('.//' + f'{{{PML_NS}}}cNvPr')
            if cNvPr is not None:
                name = cNvPr.get('name', '')
                if name != '文本框 7' and '总目录' not in name:
                    shapes_to_remove.append(sp)

    for sp in shapes_to_remove:
        spTree.remove(sp)

    # 2. 横向排开
    num_chapters = len(chapters)
    if num_chapters == 0:
         return

    # 依照历史常数布局
    chapter_box_width = 1300000 
    chapter_box_height = 461665
    chapter_gap = 120000 
    # 垂直居中
    toc_area_top = (prs.slide_height - chapter_box_height) // 2 

    slide_width = prs.slide_width
    total_width = num_chapters * chapter_box_width + (num_chapters - 1) * chapter_gap
    toc_area_left = (slide_width - total_width) // 2

    sp_id = 1000
    for ch_idx, chapter in enumerate(chapters):
        x = toc_area_left + ch_idx * (chapter_box_width + chapter_gap)
        # 添加章节 Box 元素
        chapter_sp = create_toc_chapter_box(chapter, ch_idx, x, toc_area_top, chapter_box_width, chapter_box_height, sp_id)
        spTree.append(chapter_sp)
        sp_id += 1

    print(f"  [TOC] 已为 [{num_chapters}] 个一级章节生成横向目录磁贴。")

def main():
    import argparse
    parser = argparse.ArgumentParser(description="根据大纲 JSON 生成 PPT 总目录页（第2页）")
    parser.add_argument("--input_json", default="output/mubu_parsed_structure.json", help="大纲 JSON 路径")
    parser.add_argument("--input_pptx", default="output/master_with_cover.pptx", help="第4部输出的已带封面的 PPTX 路径")
    parser.add_argument("--output_pptx", default="output/master_with_toc.pptx", help="最终带总目录的输出路径")
    args = parser.parse_args()

    input_json = Path(args.input_json)
    input_pptx = Path(args.input_pptx)
    output_pptx = Path(args.output_pptx)

    if not input_json.exists():
        print(f"❌ 大纲 JSON 文件不存在: {input_json}")
        return
    if not input_pptx.exists():
        print(f"❌ PPT 输入模板不存在: {input_pptx}")
        return

    # 1. 载入 JSON 提取一级章节
    with open(input_json, "r", encoding="utf-8") as f:
        data = json.load(f)
    chapters = data.get("chapters", [])
    if not chapters:
         print("⚠️ 警告: JSON 里面 chapters 为空。")

    # 2. 载入 PPT
    print(f"📖 打开 PPT 母本: {input_pptx}")
    prs = Presentation(input_pptx)
    if len(prs.slides) < 2:
        print("❌ PPT 没有有效的 slide (至少需要2页)，无法执行总目录替换。")
        return

    toc_slide = prs.slides[1]  # 第2页作为总目录页

    # 3. 抹除 + 生成
    print("✍️ 生成总目录视图...")
    create_toc_slide(prs, toc_slide, chapters)

    # 4. 保存
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    print(f"✅ 总目录页生成完成 -> {output_pptx}")

if __name__ == "__main__":
    main()

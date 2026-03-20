#!/usr/bin/env python3
import os
from pptx import Presentation

pptx_path = "output/final_presentation.pptx"

if not os.path.exists(pptx_path):
    print(f"File not found: {pptx_path}")
    exit()

prs = Presentation(pptx_path)
print(f"Total slides in final: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} (Layout: {slide.slide_layout.name if hasattr(slide, 'slide_layout') else 'Unknown'}) ---")
    if hasattr(slide, "shapes"):
        for j, shp in enumerate(slide.shapes):
            if shp.has_text_frame:
                txt = shp.text_frame.text.strip()
                if txt:
                    print(f"  Shape {j} [{shp.name}]: {txt[:100]}...")
            elif shp.shape_type == 13: # Picture
                print(f"  Shape {j} [Picture]: {shp.name}")

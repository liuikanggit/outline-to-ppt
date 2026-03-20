import sys
from pptx import Presentation

def inspect_ppt(path):
    print(f"正在检查 {path}...")
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"打开 PPT 失败: {e}")
        return

    print(f"总页数: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides[:10]):  # 只看前 10 页
        print(f"\n--- 第 {i+1} 页 ---")
        for shp in slide.shapes:
            if hasattr(shp, 'has_text_frame') and shp.has_text_frame:
                text = shp.text_frame.text.strip()
                if text:
                    print(f"  [文本] Name='{shp.name}' Text='{text[:30]}...'")
            elif hasattr(shp, 'image'):
                print(f"  [图片] Name='{shp.name}'")

if __name__ == '__main__':
    inspect_ppt('output/presentation.pptx')

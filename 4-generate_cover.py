#!/usr/bin/env python3
"""
第4步骤：自动生成 PPT 封面页（文字替换 + 自适应自测声明位移）
在第3步骤(create_master.py)输出的 PPTX 基础上执行。
"""

import json
import os
import math
import argparse
from pathlib import Path
from pptx import Presentation

# 如果环境无 python-pptx，请通过 venv/bin/python 运行本脚本

def main():
    parser = argparse.ArgumentParser(description="根据大纲 JSON 生成 PPT 封面（替换第一页内容）")
    parser.add_argument("--input_json", default="output/mubu_parsed_structure.json", help="大纲 JSON 路径")
    parser.add_argument("--input_pptx", default="output/master.pptx", help="第3部输出的 PPTX 母版合流文件路径")
    parser.add_argument("--output_pptx", default="output/master_with_cover.pptx", help="最终带封面的输出路径")
    args = parser.parse_args()

    input_json = Path(args.input_json)
    input_pptx = Path(args.input_pptx)
    output_pptx = Path(args.output_pptx)

    if not input_json.exists():
        print(f"❌ 大纲 JSON 文件不存在: {input_json}")
        return
    if not input_pptx.exists():
        print(f"❌ PPT 输入模板不存在: {input_pptx}")
        return

    # 1. 载入 JSON 提取课程名
    print(f"📖 读取大纲数据: {input_json}")
    with open(input_json, "r", encoding="utf-8") as f:
        data = json.load(f)
    
    # 尽可能保证两个字段均有降级容错
    course_name = data.get("courseName", "")
    course_en_name = data.get("courseEnName", data.get("chaptersEnName", ""))

    print(f"  课程中文名: {course_name}")
    print(f"  课程英文名: {course_en_name}")

    if not course_name:
        print("💡 警告: courseName 为空，将使用空白文字覆盖。")

    # 2. 载入 PPT
    print(f"📖 打开 PPT 母本: {input_pptx}")
    prs = Presentation(input_pptx)
    if not prs.slides:
        print("❌ PPT 没有有效的 slide, 无法执行封面生成。")
        return

    cover_slide = prs.slides[0]  # 第1页作为封面页

    # 3. 循环替换文字
    print("✍️ 替换封面页文字...")
    replaced_cn = False
    replaced_en = False

    for shp in cover_slide.shapes:
        if shp.has_text_frame:
            for paragraph in shp.text_frame.paragraphs:
                for run in paragraph.runs:
                    txt = run.text.strip()
                    
                    # 严格全字匹配
                    if txt == "课程名称" or txt == "XXXXX（课程名称）":
                        run.text = course_name
                        replaced_cn = True
                    elif txt == "课程名称英文" or txt == "XXXXXXXXXXX（课程名称英文）":
                        run.text = course_en_name
                        replaced_en = True
                    elif "课程名称" in txt or "XXXXX" in txt:
                        # 局部替换
                        new_txt = run.text.replace("XXXXX（课程名称）", course_name) \
                                         .replace("XXXXXXXXXXX（课程名称英文）", course_en_name) \
                                         .replace("课程名称英文", course_en_name) \
                                         .replace("课程名称", course_name) \
                                         .replace("XXXXX", "") \
                                         .replace("XXXXXXXXXXX", "")
                        run.text = new_txt
                        replaced_cn = True

    print(f"  [替换状态] 中文版替换: {'✅' if replaced_cn else '❓ 没找到占位符'}")
    print(f"  [替换状态] 英文版替换: {'✅' if replaced_en else '❓ 没找到占位符'}")

    # 4. “免责声明” 文本框自适应位移
    # 估算主标题行数 (中文汉字约14字一行, 英文按24字符一行更精确，因为英文单词换行间隙大)
    cn_lines = math.ceil(len(course_name) / 14)
    en_lines = math.ceil(len(course_en_name) / 24)
    extra = max(0, cn_lines - 1) + max(0, en_lines - 1)

    if extra > 0:
        offset = extra * 382000  # EMU 单位位移值 (382000 完美适配单行高度)
        adjusted_count = 0
        disclaimer_phrase = "未经授权严禁"
        for shp in cover_slide.shapes:
            if shp.has_text_frame and disclaimer_phrase in shp.text_frame.text:
                shp.top = shp.top + offset
                adjusted_count += 1
                print(f"  [自适应排版] 检测到课程标题多出 {extra} 行文字，声明文本框向下平移 {offset} EMU")
        if adjusted_count == 0:
             print(f"  [自适应排版] 未找到包含 '{disclaimer_phrase}' 的文本框，跳过位移。")

    # 5. 保存
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    print(f"✅ 封面生成并合成完毕 -> {output_pptx}")

if __name__ == "__main__":
    main()
